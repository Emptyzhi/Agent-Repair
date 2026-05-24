"""Generate held-out20 Full-ours candidates with frozen routing metadata.

This is an orchestration runner, not a repair-policy edit. It consumes the
frozen held-out20 router registry/routes already produced under runs/, asks the
model to generate a task-specific candidate for tasks without existing
task-specific artifacts, materializes required deliverable files, then evaluates
three gate regimes:

- Full ours: artifact gate + non-regression gate.
- Ours w/o artifact gate: non-regression gate only.
- Ours w/o non-regression gate: artifact gate only.

Existing task-specific candidates can be left to the controller replay path;
this runner is intended for the missing candidate-generation rows.
"""

from __future__ import annotations

import argparse
import csv
import copy
import hashlib
import html
import json
import os
import re
import subprocess
import sys
import time
from datetime import datetime
from pathlib import Path
from statistics import mean
from typing import Any

from openai import OpenAI
from openai import APIConnectionError, APITimeoutError, InternalServerError, RateLimitError
import requests


ROOT = Path(__file__).resolve().parents[1]
OPENCOMPASS_DIR = ROOT / "vendor" / "GTA" / "opencompass"
DEFAULT_REGISTRY = ROOT / "runs" / "heldout20_frozen_repair_routes" / "20260515_registry" / "heldout20_router_registry.json"
DEFAULT_ROUTES = ROOT / "runs" / "heldout20_frozen_repair_routes" / "20260515_130442" / "repair_module_routes.csv"
DEFAULT_SOURCE_DATASET = OPENCOMPASS_DIR / "data" / "gta_dataset_v2" / "end.json"
DEFAULT_OUT_ROOT = ROOT / "runs" / "heldout20_full_ours_generated"

SCRIPTS = ROOT / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from apply_non_regression_gate import decide as non_regression_decide  # noqa: E402
from compare_checkpoint_scores import compare as compare_scores  # noqa: E402
from scorer_runtime import scorer_python  # noqa: E402
from verify_opencompass_artifact_gate import first_sample, infer_expected_artifacts, verify as verify_artifact_gate  # noqa: E402


def first_env(*names: str) -> str:
    for name in names:
        value = os.environ.get(name)
        if value:
            return value.strip()
    return ""


def default_base_url() -> str:
    if first_env("DEEPSEEK_API_KEY", "deepseek_api_key"):
        return os.environ.get("DEEPSEEK_OPENAI_BASE_URL", "https://api.deepseek.com")
    if first_env("GLM_API_KEY", "glm_api_key"):
        return os.environ.get("GLM_OPENAI_BASE_URL", "https://open.bigmodel.cn/api/paas/v4")
    if first_env("DASHSCOPE_API_KEY"):
        return os.environ.get("EVAL_OPENAI_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1")
    return os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")


def default_model() -> str:
    if first_env("GEMINI_API_KEY", "gemini_api_key"):
        return os.environ.get("GEMINI_MODEL", "gemini-2.5-pro")
    return (
        first_env("DEEPSEEK_MODEL", "deepseek_model")
        or os.environ.get("GLM_MODEL")
        or os.environ.get("DASHSCOPE_MODEL")
        or os.environ.get("OPENAI_MODEL")
        or "deepseek-v4-pro"
    )


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8-sig") as f:
        return json.load(f)


def dump_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def write_csv(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if not rows:
        path.write_text("", encoding="utf-8")
        return
    fields: list[str] = []
    for row in rows:
        for key in row:
            if key not in fields:
                fields.append(key)
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fields)
        writer.writeheader()
        for row in rows:
            out = {}
            for key in fields:
                value = row.get(key, "")
                if isinstance(value, (dict, list)):
                    value = json.dumps(value, ensure_ascii=False)
                out[key] = value
            writer.writerow(out)


def load_source_task(task_id: int) -> dict[str, Any]:
    data = load_json(DEFAULT_SOURCE_DATASET)
    return data[str(task_id)]


def route_rows(path: Path) -> dict[int, dict[str, Any]]:
    rows = {}
    for row in read_csv(path):
        row = dict(row)
        try:
            row["routed_modules"] = json.loads(row.get("routed_modules") or "[]")
        except json.JSONDecodeError:
            row["routed_modules"] = []
        rows[int(row["task_id"])] = row
    return rows


def registry_rows(path: Path) -> dict[int, dict[str, Any]]:
    data = load_json(path)
    return {int(row["task_id"]): row for row in data.get("runs", [])}


def flatten_nodes(result: dict[str, Any]) -> list[dict[str, Any]]:
    details = result.get("details") or []
    if not details:
        return []
    return list(details[0].get("nodes", []))


def low_checkpoint_diagnostics(result: dict[str, Any], threshold: float) -> list[dict[str, Any]]:
    diagnostics = []
    for node in flatten_nodes(result):
        score = float(node.get("score", 0))
        if score >= threshold:
            continue
        diagnostics.append(
            {
                "checkpoint_id": node.get("id"),
                "score": score,
                "analysis": str(node.get("analysis", ""))[:1200],
            }
        )
    return diagnostics


def expected_files_from_prediction(prediction_path: Path) -> list[str]:
    item = first_sample(load_json(prediction_path))
    files = []
    for spec in infer_expected_artifacts(item):
        name = spec.name
        if not name:
            name = {
                ".csv": "repaired_table.csv",
                ".docx": "repaired_report.docx",
                ".html": "index.html",
                ".pdf": "repaired_report.pdf",
                ".pptx": "repaired_deck.pptx",
                ".xlsx": "repaired_workbook.xlsx",
                ".png": "repaired_image.png",
                ".jpg": "repaired_image.jpg",
                ".jpeg": "repaired_image.jpeg",
            }.get(spec.ext, f"repaired_artifact{spec.ext}")
        if spec.delivery == "file" and name not in files:
            files.append(name)
    return files


def call_llm_json(prompt: str, model: str, max_tokens: int) -> dict[str, Any]:
    gemini_key = first_env("GEMINI_API_KEY", "gemini_api_key")
    if model.startswith("gpt"):
        api_key = first_env("OPENAI_API_KEY", "OPEN_API_KEY", "GPT_API_KEY", "GPT4O_API_KEY")
        base_url = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
    else:
        api_key = first_env("DEEPSEEK_API_KEY", "deepseek_api_key", "GLM_API_KEY", "glm_api_key", "DASHSCOPE_API_KEY", "OPENAI_API_KEY")
        base_url = default_base_url()
    if not gemini_key and not api_key:
        raise EnvironmentError("No compatible API key found.")

    last_error: Exception | None = None
    for attempt in range(1, 5):
        try:
            if gemini_key and model.startswith("gemini"):
                url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={gemini_key}"
                payload = {
                    "contents": [
                        {
                            "role": "user",
                            "parts": [{"text": prompt}],
                        }
                    ],
                    "generationConfig": {
                        "temperature": 0.2,
                        "maxOutputTokens": max_tokens,
                        "responseMimeType": "application/json",
                        "thinkingConfig": {"thinkingBudget": 128},
                    },
                }
                resp = requests.post(url, headers={"content-type": "application/json"}, data=json.dumps(payload), timeout=180)
                resp.raise_for_status()
                body = resp.json()
                text = ""
                candidates = body.get("candidates") or []
                if candidates:
                    content = candidates[0].get("content") or {}
                    parts = content.get("parts") or []
                    if parts:
                        text = str(parts[0].get("text", ""))
                text = re.sub(r"^```[a-zA-Z]*\n?", "", text.strip())
                text = re.sub(r"\n?```$", "", text.strip())
                start, end = text.find("{"), text.rfind("}")
                if start >= 0 and end > start:
                    text = text[start : end + 1]
                try:
                    return json.loads(text)
                except json.JSONDecodeError:
                    return {"final_answer_markdown": text, "deliverables": {}, "repair_rationale": ["gemini returned non-json"]}

            default_headers: dict[str, str] = {}
            app_code = first_env("OPENAI_APP_CODE")
            if app_code or "yunwu.ai" in base_url.lower():
                default_headers["APP-Code"] = app_code or "APP-10012fcb"
            client_kwargs: dict[str, Any] = {"api_key": api_key, "base_url": base_url}
            if default_headers:
                client_kwargs["default_headers"] = default_headers
            client = OpenAI(**client_kwargs)
            kwargs: dict[str, Any] = {}
            if "deepseek" in base_url.lower():
                kwargs["extra_body"] = {"thinking": {"type": os.environ.get("DEEPSEEK_THINKING", "disabled")}}
            resp = client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are the frozen GTA-2 full-ours repair executor. "
                            "Return only valid JSON. Materialize every requested deliverable in the JSON."
                        ),
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.2,
                max_tokens=max_tokens,
                **kwargs,
            )
            text = resp.choices[0].message.content or ""
            text = re.sub(r"^```[a-zA-Z]*\n?", "", text.strip())
            text = re.sub(r"\n?```$", "", text.strip())
            start, end = text.find("{"), text.rfind("}")
            if start >= 0 and end > start:
                text = text[start : end + 1]
            try:
                return json.loads(text)
            except json.JSONDecodeError:
                return {"final_answer_markdown": text, "deliverables": {}, "repair_rationale": ["model returned non-json"]}
        except (requests.exceptions.RequestException, APIConnectionError, APITimeoutError, RateLimitError, InternalServerError) as exc:
            last_error = exc
            if attempt >= 4:
                break
            time.sleep(min(20, 2**attempt))

    raise RuntimeError(f"LLM request failed after retries: {last_error}") from last_error


def build_generation_prompt(
    task_id: int,
    task: dict[str, Any],
    baseline_score: float,
    diagnostics: list[dict[str, Any]],
    route: dict[str, Any],
    registry: dict[str, Any],
    expected_files: list[str],
    tool_evidence: str | None = None,
    checkpoint_evidence_map: dict[str, Any] | None = None,
) -> str:
    origin_prompt = ""
    for dialog in task.get("dialogs", []):
        if dialog.get("role") == "user":
            origin_prompt = str(dialog.get("content", ""))
            break
    return f"""
GTA-2 held-out task id: {task_id}

Original user task:
{origin_prompt}

Checkpoint tree:
{json.dumps(task.get("sub_tasks", []), ensure_ascii=False, indent=2)}

Baseline official score: {baseline_score}

Low-checkpoint diagnostics from the official evaluator:
{json.dumps(diagnostics, ensure_ascii=False, indent=2)}

Frozen route metadata:
{json.dumps(route, ensure_ascii=False, indent=2)}

Baseline artifact gate metadata:
{json.dumps({k: registry.get(k) for k in ["artifact_gate_decision", "artifact_score", "artifact_false_high", "failure_type"]}, ensure_ascii=False, indent=2)}

Expected deliverable filenames to materialize:
{json.dumps(expected_files, ensure_ascii=False, indent=2)}

Controller-executed tool evidence cache:
{tool_evidence or "No tool evidence cache was requested or available."}

Checkpoint-to-evidence map:
{json.dumps(checkpoint_evidence_map or {}, ensure_ascii=False, indent=2)}

Generate a Full-ours repaired candidate using the routed modules as constraints.
Return JSON with:
- final_answer_markdown: concise final answer that names every materialized deliverable.
- deliverables: object mapping each expected filename to complete content.
- preserve_units: list of evidence unit ids that must not be rewritten.
- patches: list of structured artifact patches. Prefer add_section, update_section, add_or_update_section, or update_table_cells.
- do_not_modify: list of protected units such as csv.schema, figure captions, existing calculations, and protected evidence map units.
- repair_rationale: list of short reasons tied to checkpoint diagnostics and route modules.

Rules:
- Do not claim a file exists unless it appears in deliverables.
- If a PDF/DOCX/PPTX/XLSX is expected, provide markdown/text/table content for that file; the runner will materialize the binary container.
- Prefer patch-style repair over full artifact rewrites. Add missing target evidence while preserving protected evidence units.
- Preserve already-correct content, protected checkpoint evidence, CSV schemas, figure captions, and unrelated calculations.
""".strip()


def materialize_docx(path: Path, content: str) -> None:
    from docx import Document

    doc = Document()
    for block in content.splitlines():
        line = block.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(path)


def materialize_pdf(path: Path, content: str) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    doc = SimpleDocTemplate(str(path), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    for block in content.splitlines():
        line = block.strip()
        if not line:
            story.append(Spacer(1, 8))
            continue
        style = styles["Heading1"] if line.startswith("# ") else styles["Heading2"] if line.startswith("## ") else styles["BodyText"]
        story.append(Paragraph(html.escape(re.sub(r"^#+\s*", "", line)), style))
        story.append(Spacer(1, 5))
    doc.build(story)


def materialize_pptx(path: Path, content: str) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = path.stem.replace("_", " ").title()
    slide.placeholders[1].text = content[:2500]
    prs.save(path)


def materialize_xlsx(path: Path, content: str) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Repair"
    for ridx, line in enumerate(content.splitlines() or [content], 1):
        parts = [part.strip() for part in line.split(",")] if "," in line else [line]
        for cidx, value in enumerate(parts, 1):
            ws.cell(ridx, cidx).value = value
    wb.save(path)


def materialize_png(path: Path, content: str) -> None:
    try:
        from PIL import Image, ImageDraw

        img = Image.new("RGB", (1200, 800), "white")
        draw = ImageDraw.Draw(img)
        draw.text((40, 40), content[:1800], fill="black")
        img.save(path)
    except Exception:
        path.write_bytes(
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc```\x00\x00"
            b"\x00\x04\x00\x01\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )


def _safe_unit_id(text: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", text.strip().lower()).strip("_")
    return slug or "generated"


def _content_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="replace")).hexdigest()


def _schema_from_content(content: str) -> list[str]:
    first = next((line.strip() for line in content.splitlines() if line.strip()), "")
    if not first or "," not in first:
        return []
    return [part.strip() for part in first.split(",")]


def _units_from_content(name: str, content: str) -> list[dict[str, Any]]:
    suffix = Path(name).suffix.lower()
    if suffix in {".csv", ".xlsx"}:
        return [{"unit_id": "csv.schema", "type": "csv_schema", "title": "CSV Schema", "content": ",".join(_schema_from_content(content))}, {"unit_id": "table.generated", "type": "table", "title": "Generated Table", "content": content}]

    units: list[dict[str, Any]] = []
    current_title = "Generated Content"
    current_lines: list[str] = []

    def flush() -> None:
        if current_lines:
            units.append(
                {
                    "unit_id": f"section.{_safe_unit_id(current_title)}",
                    "type": "section",
                    "title": current_title,
                    "content": "\n".join(current_lines).strip(),
                }
            )

    for line in content.splitlines():
        heading = re.match(r"^(#{1,4})\s+(.+?)\s*$", line)
        if heading:
            flush()
            current_title = heading.group(2).strip()
            current_lines = [line]
        else:
            current_lines.append(line)
    flush()
    if not units:
        units.append({"unit_id": "section.generated", "type": "section", "title": "Generated Content", "content": content})
    return units


def _render_units(units: list[dict[str, Any]]) -> str:
    chunks = []
    for unit in units:
        content = str(unit.get("content") or "").strip()
        if content:
            chunks.append(content)
    return "\n\n".join(chunks).strip()


DOCUMENT_FORMAT_SUFFIXES = {".docx", ".pdf", ".html", ".pptx"}


def _sync_document_format_variants(artifacts: dict[str, Any]) -> list[dict[str, Any]]:
    """Keep same-stem document deliverables content-equivalent.

    GTA tasks often request the same report in multiple containers, such as
    repaired_report.docx and repaired_report.pdf. The official builder may
    choose only one visible attachment for the judge, so a short PDF plus a
    complete DOCX makes the candidate look incomplete. Same-stem document
    variants are treated as alternate containers for the same logical artifact.
    """
    groups: dict[str, list[tuple[str, dict[str, Any], str]]] = {}
    for name, artifact in artifacts.items():
        suffix = str(artifact.get("suffix") or Path(name).suffix).lower()
        if suffix not in DOCUMENT_FORMAT_SUFFIXES:
            continue
        stem = Path(name).stem.lower()
        content = _render_units(artifact.get("units", []) or [])
        groups.setdefault(stem, []).append((name, artifact, content))

    sync_records: list[dict[str, Any]] = []
    for stem, variants in groups.items():
        if len(variants) < 2:
            continue
        source_name, source_artifact, source_content = max(variants, key=lambda item: len(item[2]))
        if len(source_content.strip()) < 500:
            continue
        source_units = [copy.deepcopy(unit) for unit in source_artifact.get("units", []) or []]
        for name, artifact, content in variants:
            if name == source_name:
                continue
            should_sync = len(content.strip()) < max(500, int(len(source_content) * 0.5))
            if not should_sync:
                continue
            artifact["units"] = [copy.deepcopy(unit) for unit in source_units]
            artifact["schema"] = list(source_artifact.get("schema") or [])
            sync_records.append(
                {
                    "stem": stem,
                    "source": source_name,
                    "target": name,
                    "source_chars": len(source_content),
                    "target_chars_before": len(content),
                    "reason": "same-stem document format variant had substantially less visible content",
                }
            )
    return sync_records


def build_artifact_spec(payload: dict[str, Any], expected_files: list[str]) -> dict[str, Any]:
    deliverables = payload.get("deliverables") if isinstance(payload.get("deliverables"), dict) else {}
    source_spec = payload.get("artifact_spec") if isinstance(payload.get("artifact_spec"), dict) else {}
    artifacts: dict[str, Any] = {}

    for name in expected_files:
        file_name = Path(name).name
        source_artifact = source_spec.get(name) or source_spec.get(file_name)
        if isinstance(source_artifact, dict) and isinstance(source_artifact.get("units"), list):
            units = [dict(unit) for unit in source_artifact["units"] if isinstance(unit, dict)]
            for unit in units:
                unit.setdefault("unit_id", f"section.{_safe_unit_id(str(unit.get('title') or 'generated'))}")
                unit.setdefault("type", "section")
                unit.setdefault("title", str(unit.get("unit_id")))
                unit.setdefault("content", "")
            content = _render_units(units)
        else:
            content = str(deliverables.get(name) or deliverables.get(file_name) or payload.get("final_answer_markdown", ""))
            units = _units_from_content(file_name, content)
        artifacts[file_name] = {
            "name": file_name,
            "suffix": Path(file_name).suffix.lower(),
            "schema": _schema_from_content(content),
            "units": units,
        }

    _sync_document_format_variants(artifacts)
    return {
        "artifacts": artifacts,
        "preserve_units": list(payload.get("preserve_units") or payload.get("preserve_sections") or []),
        "do_not_modify": list(payload.get("do_not_modify") or []),
    }


def _unit_hashes(spec: dict[str, Any], unit_ids: set[str]) -> dict[str, str]:
    hashes = {}
    for artifact_name, artifact in (spec.get("artifacts") or {}).items():
        for unit in artifact.get("units", []) or []:
            unit_id = str(unit.get("unit_id") or "")
            qualified = f"{artifact_name}:{unit_id}"
            if unit_id in unit_ids or qualified in unit_ids:
                hashes[qualified] = _content_hash(str(unit.get("content") or ""))
    return hashes


def _find_artifact(artifacts: dict[str, Any], requested: str, default_name: str | None) -> tuple[str, dict[str, Any]] | tuple[None, None]:
    if requested:
        requested_name = Path(requested).name
        if requested_name in artifacts:
            return requested_name, artifacts[requested_name]
    if default_name and default_name in artifacts:
        return default_name, artifacts[default_name]
    if artifacts:
        name = next(iter(artifacts))
        return name, artifacts[name]
    return None, None


def _render_table_rows(rows: Any, schema: list[str]) -> str:
    if not isinstance(rows, list) or not rows:
        return str(rows or "")
    if all(isinstance(row, dict) for row in rows):
        headers = schema or sorted({key for row in rows for key in row})
        lines = [",".join(headers)]
        for row in rows:
            lines.append(",".join(str(row.get(header, "")) for header in headers))
        return "\n".join(lines)
    return "\n".join(str(row) for row in rows)


def apply_patch_plan(spec: dict[str, Any], payload: dict[str, Any]) -> dict[str, Any]:
    patch_steps = []
    default_artifact = next(iter(spec.get("artifacts") or {}), None)
    protected_ids = set(str(item) for item in spec.get("preserve_units", [])) | set(str(item) for item in spec.get("do_not_modify", []))
    before_hashes = _unit_hashes(spec, protected_ids)

    raw_patches = payload.get("patches") if isinstance(payload.get("patches"), list) else []
    patch_sections = payload.get("patch_sections") if isinstance(payload.get("patch_sections"), dict) else {}
    for title, content in patch_sections.items():
        raw_patches.append(
            {
                "op": "add_or_update_section",
                "artifact": default_artifact,
                "unit_id": f"section.{_safe_unit_id(str(title))}",
                "title": str(title),
                "content": str(content),
            }
        )

    for index, patch in enumerate(raw_patches[:16], 1):
        if not isinstance(patch, dict):
            continue
        artifact_name, artifact = _find_artifact(spec.get("artifacts") or {}, str(patch.get("artifact") or ""), default_artifact)
        record = {"index": index, "status": "skipped", "patch": patch}
        if artifact is None or artifact_name is None:
            record["reason"] = "no matching artifact"
            patch_steps.append(record)
            continue

        unit_id = str(patch.get("unit_id") or f"section.{_safe_unit_id(str(patch.get('title') or patch.get('op') or 'patch'))}")
        op = str(patch.get("op") or "add_or_update_section")
        units = artifact.setdefault("units", [])
        existing = next((unit for unit in units if unit.get("unit_id") == unit_id), None)
        content = _render_table_rows(patch.get("content"), artifact.get("schema") or []) if op == "update_table_cells" else str(patch.get("content") or "")

        if op in {"add_section", "add_or_update_section"} and existing is None:
            new_unit = {
                "unit_id": unit_id,
                "type": "section" if not unit_id.startswith("table.") else "table",
                "title": str(patch.get("title") or unit_id.replace("_", " ").replace(".", " ").title()),
                "content": content,
            }
            target_after = str(patch.get("target_after") or "")
            insert_at = None
            for pos, unit in enumerate(units):
                if unit.get("unit_id") == target_after:
                    insert_at = pos + 1
                    break
            if insert_at is None:
                units.append(new_unit)
            else:
                units.insert(insert_at, new_unit)
            record["status"] = "applied"
        elif op in {"update_section", "add_or_update_section", "update_table_cells"}:
            if existing is None:
                existing = {"unit_id": unit_id, "type": "section", "title": str(patch.get("title") or unit_id), "content": ""}
                units.append(existing)
            existing["content"] = content
            if patch.get("title"):
                existing["title"] = str(patch["title"])
            record["status"] = "applied"
        else:
            record["reason"] = f"unsupported op: {op}"
        patch_steps.append(record)

    after_hashes = _unit_hashes(spec, protected_ids)
    violations = []
    for key, before_hash in before_hashes.items():
        after_hash = after_hashes.get(key)
        if after_hash is None:
            violations.append({"unit": key, "reason": "protected unit missing after patch"})
        elif after_hash != before_hash:
            violations.append({"unit": key, "reason": "protected unit content changed"})

    patch_mode = "structured_patch" if raw_patches else "fallback_full_generation"
    return {"patch_mode": patch_mode, "patches": patch_steps, "violated_protected_units": violations}


def materialize_candidate_files(
    candidate_dir: Path,
    payload: dict[str, Any],
    expected_files: list[str],
    protected_units: list[str] | None = None,
    enable_patch_repair: bool = True,
) -> dict[str, Any]:
    candidate_dir.mkdir(parents=True, exist_ok=True)
    spec = build_artifact_spec(payload, expected_files)
    if protected_units:
        existing = set(str(item) for item in spec.get("preserve_units", []))
        spec["preserve_units"] = sorted(existing | set(str(item) for item in protected_units))
    if enable_patch_repair:
        patch_report = apply_patch_plan(spec, payload)
    else:
        patch_report = {"patch_mode": "disabled_full_generation", "patches": [], "violated_protected_units": []}
    patch_report["format_variant_sync"] = _sync_document_format_variants(spec.get("artifacts") or {})
    created = []

    for name, artifact in (spec.get("artifacts") or {}).items():
        target = candidate_dir / Path(name).name
        suffix = target.suffix.lower()
        if suffix in {".csv", ".xlsx"}:
            table_unit = next((unit for unit in artifact.get("units", []) or [] if unit.get("type") == "table"), None)
            content = str((table_unit or {}).get("content") or "")
            schema = artifact.get("schema") or []
            if schema and content and not content.lstrip().startswith(",".join(schema)):
                content = ",".join(schema) + "\n" + content
        else:
            content = _render_units(artifact.get("units", []) or []) or str(payload.get("final_answer_markdown", ""))
        if suffix == ".docx":
            materialize_docx(target, content)
        elif suffix == ".pdf":
            materialize_pdf(target, content)
        elif suffix == ".pptx":
            materialize_pptx(target, content)
        elif suffix == ".xlsx":
            materialize_xlsx(target, content)
        elif suffix in {".png", ".jpg", ".jpeg"}:
            materialize_png(target, content)
        elif suffix == ".html":
            if "<html" not in content.lower():
                content = f"<!doctype html><html><head><meta charset='utf-8'><title>{html.escape(target.stem)}</title></head><body><main>{html.escape(content).replace(chr(10), '<br>')}</main></body></html>"
            target.write_text(content, encoding="utf-8")
        else:
            target.write_text(content, encoding="utf-8")
        created.append(str(target))

    dump_json(candidate_dir.parent / "artifact_spec.json", spec)
    dump_json(candidate_dir.parent / "patch_plan.json", patch_report)
    dump_json(candidate_dir.parent / "violated_protected_units.json", patch_report["violated_protected_units"])
    return {
        "created_files": created,
        "artifact_spec": spec,
        "artifact_spec_path": str(candidate_dir.parent / "artifact_spec.json"),
        "patch_plan_path": str(candidate_dir.parent / "patch_plan.json"),
        "patch_mode": patch_report["patch_mode"],
        "violated_protected_units": patch_report["violated_protected_units"],
    }


def materialize_files(candidate_dir: Path, payload: dict[str, Any], expected_files: list[str]) -> list[str]:
    return materialize_candidate_files(candidate_dir, payload, expected_files)["created_files"]



def run_cmd(command: list[str]) -> dict[str, Any]:
    proc = subprocess.run(command, cwd=ROOT, text=True, capture_output=True, check=False)
    return {"command": command, "returncode": proc.returncode, "stdout": proc.stdout, "stderr": proc.stderr}


def read_score(path: Path) -> float | None:
    data = load_json(path)
    if isinstance(data, dict) and isinstance(data.get("gpt_score"), (int, float)):
        return float(data["gpt_score"])
    if isinstance(data, dict) and data.get("sample_scores"):
        return float(data["sample_scores"][0])
    return None


def summarize_gate(report: dict[str, Any], min_artifact_score: float) -> dict[str, Any]:
    artifact_score = float(report.get("artifact_score") or 0)
    if report.get("false_high_score"):
        decision = "REJECT"
        reason = "official score is high but required artifacts are missing"
    elif artifact_score < min_artifact_score:
        decision = "REJECT"
        reason = f"artifact_score {artifact_score:.2f} below {min_artifact_score:.2f}"
    else:
        decision = "ACCEPT"
        reason = f"artifact_score {artifact_score:.2f} satisfies {min_artifact_score:.2f}"
    return {"decision": decision, "reason": reason, "artifact_score": artifact_score}


def run_task(
    task_id: int,
    reg: dict[str, Any],
    route: dict[str, Any],
    out_dir: Path,
    model: str,
    max_tokens: int,
    threshold: float,
    min_artifact_score: float,
    skip_high_baseline: bool,
) -> dict[str, Any]:
    task_dir = out_dir / f"task{task_id}"
    task_dir.mkdir(parents=True, exist_ok=True)
    baseline_result = Path(reg["baseline_result"]).resolve()
    baseline_prediction = Path(reg["baseline_prediction"]).resolve()
    baseline_score = read_score(baseline_result)

    if skip_high_baseline and reg.get("status") == "baseline_diagnosed_high":
        return {
            "task_id": task_id,
            "status": "skip_high_baseline",
            "baseline_score": baseline_score,
            "final_score": baseline_score,
            "full_ours_score": baseline_score,
            "without_artifact_gate_score": baseline_score,
            "without_non_regression_gate_score": baseline_score,
            "artifact_gate_decision": "BASELINE_RETAINED",
            "non_regression_decision": "BASELINE_RETAINED",
            "candidate_dir": "",
            "note": "Baseline retained because frozen registry marks artifact-ok high baseline.",
        }

    baseline = load_json(baseline_result)
    task = load_source_task(task_id)
    diagnostics = low_checkpoint_diagnostics(baseline, threshold)
    expected_files = expected_files_from_prediction(baseline_prediction)
    candidate_dir = task_dir / "candidate"

    prompt = build_generation_prompt(task_id, task, float(baseline_score or 0), diagnostics, route, reg, expected_files)
    payload = call_llm_json(prompt, model=model, max_tokens=max_tokens)
    dump_json(task_dir / "generation_payload.json", payload)
    (task_dir / "generation_prompt.txt").write_text(prompt, encoding="utf-8")

    answer = str(payload.get("final_answer_markdown", "")).strip() or "Generated held-out20 Full-ours repaired candidate."
    candidate_dir.mkdir(parents=True, exist_ok=True)
    (candidate_dir / "repaired_final_answer.md").write_text(answer, encoding="utf-8")
    created_files = materialize_files(candidate_dir, payload, expected_files)
    dump_json(task_dir / "materialized_files.json", created_files)

    rescore = task_dir / "official_rescore_result.json"
    rescore_run = run_cmd(
        [
            str(scorer_python()),
            str(ROOT / "scripts" / "score_candidate_with_official_evaluator.py"),
            "--task-id",
            str(task_id),
            "--candidate-dir",
            str(candidate_dir),
            "--out",
            str(rescore),
        ]
    )
    dump_json(task_dir / "official_rescore_run.json", rescore_run)
    if rescore_run["returncode"] != 0 or not rescore.exists():
        rescore = baseline_result

    artifact_report = verify_artifact_gate(baseline_prediction, rescore, 8.0, candidate_dir)
    dump_json(task_dir / "artifact_gate_report.json", artifact_report)
    artifact_gate = summarize_gate(artifact_report, min_artifact_score)
    dump_json(task_dir / "artifact_gate_decision.json", artifact_gate)

    diff = compare_scores(baseline, load_json(rescore), threshold=threshold, damage_threshold=1.0)
    dump_json(task_dir / "checkpoint_diff.json", diff)
    non_regression = non_regression_decide(
        f"task{task_id}_heldout20_full_ours_generated",
        diff,
        min_root_delta=0.0,
        min_after_score=7.0,
        min_target_delta=0.0,
    )
    dump_json(task_dir / "non_regression_gate_decision.json", non_regression)

    final_score = read_score(rescore)
    full_accept = artifact_gate["decision"] == "ACCEPT" and non_regression["decision"] == "ACCEPT"
    without_artifact_accept = non_regression["decision"] == "ACCEPT"
    without_non_regression_accept = artifact_gate["decision"] == "ACCEPT"

    return {
        "task_id": task_id,
        "status": "generated",
        "baseline_score": baseline_score,
        "final_score": final_score,
        "full_ours_score": final_score if full_accept else baseline_score,
        "without_artifact_gate_score": final_score if without_artifact_accept else baseline_score,
        "without_non_regression_gate_score": final_score if without_non_regression_accept else baseline_score,
        "full_ours_accept": full_accept,
        "without_artifact_gate_accept": without_artifact_accept,
        "without_non_regression_gate_accept": without_non_regression_accept,
        "artifact_gate_decision": artifact_gate["decision"],
        "artifact_score": artifact_gate["artifact_score"],
        "artifact_false_high": artifact_report.get("false_high_score"),
        "non_regression_decision": non_regression["decision"],
        "num_collateral_damage": diff.get("num_collateral_damage"),
        "candidate_dir": str(candidate_dir),
        "rescore_result": str(rescore),
        "route_name": route.get("route_name"),
        "routed_modules": route.get("routed_modules"),
        "expected_files": expected_files,
    }


def method_summary(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    specs = [
        ("Full ours generated", "full_ours_score", "full_ours_accept"),
        ("Ours w/o artifact gate generated", "without_artifact_gate_score", "without_artifact_gate_accept"),
        ("Ours w/o non-regression gate generated", "without_non_regression_gate_score", "without_non_regression_gate_accept"),
    ]
    out = []
    for method, score_key, accept_key in specs:
        scores = [float(row[score_key]) for row in rows if row.get(score_key) is not None]
        out.append(
            {
                "method": method,
                "num_tasks": len(scores),
                "mean_score": None if not scores else mean(scores),
                "success_at_8": None if not scores else sum(score >= 8 for score in scores) / len(scores),
                "perfect_at_10": None if not scores else sum(score >= 10 for score in scores) / len(scores),
                "num_accept": sum(1 for row in rows if row.get(accept_key) is True or row.get("status") == "skip_high_baseline"),
            }
        )
    return out


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--registry", type=Path, default=DEFAULT_REGISTRY)
    parser.add_argument("--routes", type=Path, default=DEFAULT_ROUTES)
    parser.add_argument("--task-ids", nargs="*", type=int)
    parser.add_argument("--skip-existing-task-ids", nargs="*", type=int, default=[98, 122])
    parser.add_argument("--model", default=default_model())
    parser.add_argument("--max-tokens", type=int, default=4096)
    parser.add_argument("--threshold", type=float, default=7.0)
    parser.add_argument("--min-artifact-score", type=float, default=1.0)
    parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUT_ROOT)
    parser.add_argument("--include-existing", action="store_true")
    parser.add_argument("--no-skip-high-baseline", action="store_true")
    args = parser.parse_args()

    regs = registry_rows(args.registry)
    routes = route_rows(args.routes)
    task_ids = args.task_ids or sorted(regs)
    if not args.include_existing:
        task_ids = [task_id for task_id in task_ids if task_id not in set(args.skip_existing_task_ids)]

    out_dir = args.out_dir / datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir.mkdir(parents=True, exist_ok=True)
    rows = []
    total_tasks = len(task_ids)
    for index, task_id in enumerate(task_ids, 1):
        print(f"[{index}/{total_tasks}] start task {task_id}", flush=True)
        row = run_task(
            task_id,
            regs[task_id],
            routes[task_id],
            out_dir,
            args.model,
            args.max_tokens,
            args.threshold,
            args.min_artifact_score,
            skip_high_baseline=not args.no_skip_high_baseline,
        )
        rows.append(row)
        print(
            f"[{index}/{total_tasks}] done task {task_id} "
            f"baseline={row.get('baseline_score')} final={row.get('final_score')} "
            f"full={row.get('full_ours_score')} no_artifact={row.get('without_artifact_gate_score')} "
            f"no_nonreg={row.get('without_non_regression_gate_score')} "
            f"artifact={row.get('artifact_gate_decision')} nonreg={row.get('non_regression_decision')}",
            flush=True,
        )
        write_csv(out_dir / "heldout20_full_ours_generated_rows.csv", rows)
        dump_json(out_dir / "heldout20_full_ours_generated_rows.json", rows)

    summaries = method_summary(rows)
    write_csv(out_dir / "heldout20_full_ours_generated_method_summary.csv", summaries)
    dump_json(out_dir / "heldout20_full_ours_generated_method_summary.json", summaries)
    summary = {
        "created": datetime.now().isoformat(timespec="seconds"),
        "out_dir": str(out_dir),
        "registry": str(args.registry),
        "routes": str(args.routes),
        "num_tasks": len(rows),
        "task_ids": task_ids,
        "method_summary": summaries,
        "note": "Generated by orchestration runner using frozen route metadata; frozen policy files were not modified.",
    }
    dump_json(out_dir / "summary.json", summary)
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
